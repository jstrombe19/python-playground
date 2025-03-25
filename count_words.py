from pptx import Presentation
import sys

def count_words_in_slide(slide):
    word_count = 0
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            word_count += len(shape.text.split())
    return word_count

def count_words_per_slide(pptx_file):
    prs = Presentation(pptx_file)
    slide_word_counts = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        word_count = count_words_in_slide(slide)
        slide_word_counts.append((slide_number, word_count))
    
    return slide_word_counts

def main(pptx_file):
    slide_word_counts = count_words_per_slide(pptx_file)
    for slide_number, word_count in slide_word_counts:
        print(f"Slide {slide_number}: {word_count} words")

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python script.py presentation.pptx")
    else:
        pptx_file = sys.argv[1]
        main(pptx_file)
