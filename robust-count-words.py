from pptx import Presentation
import sys

def count_words_in_shape(shape):
    word_count = 0
    if hasattr(shape, "text"):
        word_count += len(shape.text.split())
    elif shape.shape_type == 19:  # GROUP_SHAPE
        for sub_shape in shape.shapes:
            word_count += count_words_in_shape(sub_shape)
    return word_count

def count_words_in_table(table):
    word_count = 0
    for row in table.rows:
        for cell in row.cells:
            if cell.text:
                word_count += len(cell.text.split())
    return word_count

def count_words_in_chart(chart):
    word_count = 0
    # Count words in chart title
    if chart.has_title:
        word_count += len(chart.chart_title.text_frame.text.split())
    # Count words in data labels, if any
    for series in chart.series:
        if series.has_data_labels:
            for data_label in series.data_labels:
                if data_label.has_text_frame:
                    word_count += len(data_label.text_frame.text.split())
    return word_count

def count_words_in_slide(slide):
    word_count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            word_count += count_words_in_shape(shape)
        elif shape.has_table:
            word_count += count_words_in_table(shape.table)
        elif shape.has_chart:
            word_count += count_words_in_chart(shape.chart)
    return word_count

def count_words_per_slide(pptx_file):
    prs = Presentation(pptx_file)
    slide_word_counts = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        word_count = count_words_in_slide(slide)
        slide_word_counts.append((slide_number, word_count))
    
    return slide_word_counts

def main(pptx_file):
    slide_word_counts = count_words_per_slide(pptx_file)
    for slide_number, word_count in slide_word_counts:
        print(f"Slide {slide_number}: {word_count} words")

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python script.py presentation.pptx")
    else:
        pptx_file = sys.argv[1]
        main(pptx_file)

